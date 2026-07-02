#!/usr/bin/env python3
"""Extract text and images from PPTX files in a cross-platform way.

Overview:
    - Extract text with ``python-pptx``, one entry per slide separated by ``---``.
    - Extract embedded raster images with ``python-pptx`` (PNG/JPEG blobs saved
      as PNG).
    - Avoid OS-specific tools. No LibreOffice or system binaries required.

Usage:
    Basic run over all PPTX files in ``inputs/content-to-ingest/``::

        python tools/extract_pptx_assets.py

    Specify image background mode and output location::

        python tools/extract_pptx_assets.py \
            --background transparent \
            --output-dir outputs/pptx-processing

    Intentionally overwrite existing matching output directories::

        python tools/extract_pptx_assets.py --overwrite

Output layout:
    For each ``inputs/content-to-ingest/<name>.pptx``, this script creates:

    - ``outputs/pptx-processing/<name>/text.md``
    - ``outputs/pptx-processing/<name>/resources/*.png``

    Text is organised per slide with slide numbers and titles where present.
    Each slide section is separated by ``---`` to match the PDF extractor format.

    By default, name collisions create suffixed directories (for example,
    ``slides-2``) to avoid overwriting previous outputs.

Known limitations:
    - Vector shapes, SmartArt, charts, and drawn shapes are not rasterised.
      Their text content (if any) is extracted; their visual appearance is not.
      A future ``--render-vector-slides`` option may address this via an
      optional LibreOffice fallback.
    - Slide background images are not extracted; only shapes with image fills
      and picture placeholders are captured.

Requirements:
    Install dependencies from ``requirements.txt`` or directly::

        python -m pip install python-pptx Pillow
"""

from __future__ import annotations

import argparse
import io
import json
import os
import sys
from dataclasses import dataclass
from pathlib import Path

from dotenv import load_dotenv
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def _find_repo_root() -> Path:
    for candidate in Path(__file__).resolve().parents:
        if (candidate / "README.md").exists() and (candidate / ".github").exists():
            return candidate
    raise RuntimeError("Unable to determine repository root from script location.")


ROOT_DIR = _find_repo_root()
load_dotenv(ROOT_DIR / ".env")


def _resolve_dir_from_env(var_name: str, fallback: str) -> Path:
    value = os.getenv(var_name, fallback).strip()
    path = Path(value)
    if not path.is_absolute():
        path = ROOT_DIR / path
    return path


DEFAULT_INPUT_DIR = (
    _resolve_dir_from_env("CONTENT_INGESTER_INPUTS_DIR", "inputs") / "content-to-ingest"
)
DEFAULT_OUTPUT_DIR = (
    _resolve_dir_from_env("CONTENT_INGESTER_OUTPUTS_DIR", "outputs") / "pptx-processing"
)


@dataclass
class PptxExtractionResult:
    """Result summary for a single processed PPTX file.

    Attributes:
        pptx_path: Input PPTX path.
        output_dir: Output directory for this file.
        text_path: Path to extracted text file.
        embedded_images: Number of embedded images exported.
        skipped_vector_shapes: Number of shapes with no raster image (vector /
            chart / SmartArt) that were skipped for image export.
    """

    pptx_path: str
    output_dir: str
    text_path: str
    embedded_images: int
    skipped_vector_shapes: int


def _resolve_output_dir(output_root: Path, pptx_path: Path, overwrite: bool) -> Path:
    """Resolve a non-colliding output directory for a PPTX file.

    Args:
        output_root: Root directory for all extraction outputs.
        pptx_path: Source PPTX path.
        overwrite: Whether existing output for the same stem may be reused.

    Returns:
        A directory path that is safe to write to.
    """
    base_dir = output_root / pptx_path.stem
    if overwrite or not base_dir.exists():
        return base_dir

    suffix = 2
    while True:
        candidate = output_root / f"{pptx_path.stem}-{suffix}"
        if not candidate.exists():
            return candidate
        suffix += 1


def _save_image_blob(blob: bytes, output_path: Path, background: str) -> None:
    """Convert an image blob to PNG and save it with the requested background policy.

    Uses Pillow to handle any source format (PNG, JPEG, BMP, etc.) and to
    apply the same transparent / white / opaque background policy as the PDF
    extractor.

    Args:
        blob: Raw image bytes from a pptx image part.
        output_path: Destination ``.png`` file path.
        background: Background mode: ``transparent``, ``white``, or ``opaque``.

    Returns:
        None.
    """
    img = Image.open(io.BytesIO(blob))

    if background == "transparent":
        # Ensure image has an alpha channel so transparency is preserved.
        if img.mode not in ("RGBA", "LA"):
            img = img.convert("RGBA")
        img.save(output_path, format="PNG")

    elif background in ("white", "opaque"):
        # Flatten any alpha onto a white background, then save without alpha.
        if img.mode in ("RGBA", "LA"):
            if img.mode == "LA":
                img = img.convert("RGBA")
            bg = Image.new("RGB", img.size, (255, 255, 255))
            # Use the alpha channel as a paste mask to composite onto white.
            bg.paste(img, mask=img.split()[3])
            img = bg
        else:
            img = img.convert("RGB")
        img.save(output_path, format="PNG")

    else:
        img.save(output_path, format="PNG")


def _slide_title(slide) -> str | None:  # type: ignore[no-untyped-def]
    """Extract the title string from a slide, if present.

    Args:
        slide: A ``pptx.slide.Slide`` object.

    Returns:
        The title text, or ``None`` if no title placeholder is found.
    """
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()
        return title if title else None
    return None


def _extract_slide_text(slide, slide_number: int) -> str:
    """Extract all text from a single slide in reading order.

    Text is gathered from all shapes with text frames, skipping shapes whose
    sole content is the slide title (already emitted in the header line).

    Args:
        slide: A ``pptx.slide.Slide`` object.
        slide_number: 1-based slide index used in the header line.

    Returns:
        Formatted text block for this slide.
    """
    title = _slide_title(slide)
    header = f"## Slide {slide_number}"
    if title:
        header += f": {title}"

    lines: list[str] = [header]
    title_shape = slide.shapes.title

    for shape in slide.shapes:
        # Skip the title shape — already in the header.
        if shape == title_shape:
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)

    return "\n".join(lines)


def _extract_text(prs: Presentation) -> str:
    """Extract per-slide text from a Presentation.

    Args:
        prs: Open ``pptx.Presentation`` object.

    Returns:
        UTF-8 text with slide separators suitable for saving to ``text.md``.
    """
    slides: list[str] = []
    for slide_index, slide in enumerate(prs.slides):
        slide_text = _extract_slide_text(slide, slide_index + 1)
        slides.append(slide_text)

    return "\n\n---\n\n".join(slides).strip() + "\n"


def extract_pptx(
    pptx_path: Path,
    output_root: Path,
    background: str,
    overwrite: bool,
) -> PptxExtractionResult:
    """Extract text and image assets from one PPTX file.

    Args:
        pptx_path: Source PPTX path.
        output_root: Root folder where this file's output folder is created.
        background: Image background mode (``transparent``, ``white``, or
            ``opaque``).
        overwrite: Whether existing output directory for the same stem should
            be reused.

    Returns:
        A ``PptxExtractionResult`` describing generated artifacts.
    """
    output_dir = _resolve_output_dir(output_root, pptx_path, overwrite)
    output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation(pptx_path.as_posix())

    # 1) Text extraction.
    text_path = output_dir / "text.md"
    text_path.write_text(_extract_text(prs), encoding="utf-8")

    # 2) Image extraction.
    resources_dir = output_dir / "resources"
    resources_dir.mkdir(parents=True, exist_ok=True)

    embedded_images = 0
    skipped_vector_shapes = 0

    for slide_index, slide in enumerate(prs.slides):
        img_index = 0
        for shape in slide.shapes:
            # MSO_SHAPE_TYPE.PICTURE covers standard picture shapes and
            # picture placeholders filled with an image.
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                output_name = (
                    f"s{slide_index + 1:03d}_img{img_index + 1:03d}.png"
                )
                _save_image_blob(image.blob, resources_dir / output_name, background)
                embedded_images += 1
                img_index += 1
            elif shape.shape_type in (
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.GROUP,
                # Charts and SmartArt do not have a simple raster blob.
                # They are counted as skipped vector shapes.
                MSO_SHAPE_TYPE.CHART,
                MSO_SHAPE_TYPE.IGX_GRAPHIC,  # SmartArt in python-pptx 1.x
                MSO_SHAPE_TYPE.CANVAS,
                MSO_SHAPE_TYPE.DIAGRAM,
            ):
                # Only count as skipped if the shape is visually significant
                # (not a text-only auto shape).
                if not shape.has_text_frame:
                    skipped_vector_shapes += 1

    return PptxExtractionResult(
        pptx_path=pptx_path.as_posix(),
        output_dir=output_dir.as_posix(),
        text_path=text_path.as_posix(),
        embedded_images=embedded_images,
        skipped_vector_shapes=skipped_vector_shapes,
    )


def parse_args() -> argparse.Namespace:
    """Parse command-line arguments.

    Returns:
        Parsed CLI namespace.
    """
    parser = argparse.ArgumentParser(
        description="Extract text and images from PPTX files in a directory.",
    )
    parser.add_argument(
        "--input-dir",
        type=Path,
        default=DEFAULT_INPUT_DIR,
        help=(
            "Directory containing PPTX files "
            "(default: <inputs-dir>/content-to-ingest)."
        ),
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=DEFAULT_OUTPUT_DIR,
        help=(
            "Directory for extracted assets "
            "(default: <outputs-dir>/pptx-processing)."
        ),
    )
    parser.add_argument(
        "--background",
        choices=["transparent", "white", "opaque"],
        default="transparent",
        help="Background mode for exported images (default: transparent).",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help=(
            "Reuse existing output directory names (may overwrite files). "
            "By default, the script creates a unique suffixed directory "
            "when name collisions occur."
        ),
    )
    return parser.parse_args()


def main() -> int:
    """Run the CLI workflow.

    Returns:
        Process exit code. ``0`` means success, ``1`` means a recoverable
        input/configuration error.
    """
    args = parse_args()

    if not args.input_dir.exists():
        print(f"Input directory not found: {args.input_dir}", file=sys.stderr)
        return 1

    pptx_files = sorted(args.input_dir.glob("*.pptx"))
    if not pptx_files:
        print(f"No PPTX files found in {args.input_dir}", file=sys.stderr)
        return 1

    args.output_dir.mkdir(parents=True, exist_ok=True)

    all_results: list[PptxExtractionResult] = []
    for pptx_path in pptx_files:
        extraction = extract_pptx(
            pptx_path=pptx_path,
            output_root=args.output_dir,
            background=args.background,
            overwrite=args.overwrite,
        )
        all_results.append(extraction)

    summary = {
        "input_dir": args.input_dir.as_posix(),
        "output_dir": args.output_dir.as_posix(),
        "background": args.background,
        "pptx_count": len(all_results),
        "results": [r.__dict__ for r in all_results],
    }
    print(json.dumps(summary, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
